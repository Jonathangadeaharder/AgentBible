#!/usr/bin/env python3
"""Build deck.pptx from rendered slide PNGs + presenter notes.

Called by render_deck.sh after the PNG export; can also be run standalone:

    python3 build_pptx.py <deck-dir>

Reads:
  <deck-dir>/renders/slide-*.png   (sorted numerically)
  <deck-dir>/slides.md             (notes = last <!-- ... --> block per slide)

Writes:
  <deck-dir>/deck.pptx             (full-bleed 16:9 slides, notes embedded)

Why PPTX: it is the only export format Keynote imports *with* speaker notes.
Requires python-pptx (pip install python-pptx [--break-system-packages]).
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

SLIDE_WIDTH = Emu(9144000)   # 10 in
SLIDE_HEIGHT = Emu(5143500)  # 5.625 in (16:9)


def collect_pngs(render_dir: Path) -> list[Path]:
    return sorted(
        render_dir.glob("slide-*.png"),
        key=lambda p: int(re.search(r"(\d+)", p.stem).group(1)),  # type: ignore[union-attr]
    )


def _is_frontmatter(chunk: str) -> bool:
    """True if chunk is a Slidev per-slide frontmatter block (YAML only)."""
    for line in chunk.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if re.match(r"^[a-zA-Z_][a-zA-Z0-9_-]*\s*:", stripped):
            continue
        if re.match(r"^\s+\S", line):  # indented YAML continuation
            continue
        return False
    return True


def extract_notes(md_path: Path) -> list[str]:
    """One note string per slide (empty if the slide has no notes)."""
    text = md_path.read_text(encoding="utf-8")
    raw_chunks = re.split(r"\n---[ \t]*\n", "\n" + text + "\n")
    chunks = [c.strip() for c in raw_chunks if c.strip()]

    # Merge frontmatter-only chunks into the slide content that follows them.
    slides: list[str] = []
    pending_fm: str | None = None
    for chunk in chunks:
        if _is_frontmatter(chunk):
            pending_fm = chunk
        else:
            if pending_fm is not None:
                chunk = pending_fm + "\n" + chunk
                pending_fm = None
            slides.append(chunk)
    if pending_fm is not None:
        slides.append(pending_fm)

    notes: list[str] = []
    for slide_content in slides:
        matches = list(re.finditer(r"<!--(.*?)-->", slide_content, re.DOTALL))
        if matches:
            raw = matches[-1].group(1).strip()
            raw = re.sub(r"^Presenter\s+Notes\s*:\s*\n?", "", raw, flags=re.IGNORECASE)
            notes.append(raw.strip())
        else:
            notes.append("")
    return notes


def build(deck_dir: Path) -> None:
    renders = deck_dir / "renders"
    slides_md = deck_dir / "slides.md"
    output = deck_dir / "deck.pptx"

    pngs = collect_pngs(renders)
    if not pngs:
        sys.exit(f"!! no slide-*.png in {renders} — run the PNG export first")

    notes = extract_notes(slides_md) if slides_md.exists() else []
    if len(notes) != len(pngs):
        print(
            f"!! note count ({len(notes)}) != PNG count ({len(pngs)}) — "
            "likely a separator bug in slides.md; notes matched by index",
            file=sys.stderr,
        )

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    for i, png in enumerate(pngs):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        note = notes[i] if i < len(notes) else ""
        if note:
            slide.notes_slide.notes_text_frame.text = note

    prs.save(str(output))
    print(f"==> wrote {output} ({len(pngs)} slides, notes on "
          f"{sum(1 for n in notes if n)} of them)")


if __name__ == "__main__":
    if len(sys.argv) != 2:
        sys.exit("usage: build_pptx.py <deck-dir>")
    build(Path(sys.argv[1]).resolve())
